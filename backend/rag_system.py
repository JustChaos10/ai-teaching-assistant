import os
import glob
import tempfile
import re
import time
import shutil
from typing import List, Optional, Dict, Any
from dotenv import load_dotenv
import pypdf
import docx2txt
from pptx import Presentation
from langchain_groq import ChatGroq
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

def extract_text_from_pdf(file_path: str) -> str:
    reader = pypdf.PdfReader(file_path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)

def extract_text_from_docx(file_path: str) -> str:
    return docx2txt.process(file_path)

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return "\n".join(lines)

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def load_document_from_path(file_path: str) -> List[Dict[str, Any]]:
    ext = os.path.splitext(file_path)[1].lower()
    docs = []
    filename = os.path.basename(file_path)
    
    try:
        if ext == ".pdf":
            content = extract_text_from_pdf(file_path)
        elif ext == ".docx":
            content = extract_text_from_docx(file_path)
        elif ext == ".pptx":
            content = extract_text_from_pptx(file_path)
        elif ext == ".txt":
            content = extract_text_from_txt(file_path)
        else:
            return docs
            
        docs.append({
            "content": content,
            "filename": filename,
            "subject": "math" if "aejm" in filename.lower() else "reading" if "aemr" in filename.lower() else "general"
        })
    except Exception as e:
        pass
    
    return docs

def smart_chunk_documents(documents: List[Dict[str, Any]]) -> List[Document]:
    chunks = []
    
    for doc_data in documents:
        content = doc_data["content"]
        filename = doc_data["filename"]
        subject = doc_data["subject"]
        
        sections = re.split(r'\n(?=Chapter|Lesson|Unit|Exercise|Activity|\d+\.)', content)
        
        for section_idx, section in enumerate(sections):
            if len(section.strip()) < 50:
                continue
                
            if len(section) > 1000:
                words = section.split()
                for i in range(0, len(words), 300):
                    chunk_words = words[i:i+400]
                    chunk_text = " ".join(chunk_words)
                    
                    if len(chunk_text.strip()) > 30:
                        doc_obj = Document(
                            page_content=chunk_text,
                            metadata={
                                "source": filename,
                                "subject": subject,
                                "section": section_idx,
                                "chunk_type": "content"
                            }
                        )
                        chunks.append(doc_obj)
            else:
                doc_obj = Document(
                    page_content=section,
                    metadata={
                        "source": filename,
                        "subject": subject,
                        "section": section_idx,
                        "chunk_type": "section"
                    }
                )
                chunks.append(doc_obj)
    
    return chunks

class RAGSystem:
    def __init__(self, doc_folder: str = "./docs", index_folder: str = "./indexes"):
        self.doc_folder = doc_folder
        self.index_folder = index_folder
        self.faiss_index_path = os.path.join(self.index_folder, "faiss_index")
        
        self.conversation_history = []
        self.current_subject = None
        self.learning_progress = {}
        
        os.makedirs(self.doc_folder, exist_ok=True)
        os.makedirs(self.index_folder, exist_ok=True)
        
        self.embeddings, self.embeddings_available = self._load_embeddings_with_retry()
        
        groq_api_key = os.getenv("GROQ_API_KEY")
        if not groq_api_key:
            print("[WARNING] GROQ_API_KEY not found in environment variables")
            print(f"Looking for .env file at: {os.path.join(os.path.dirname(__file__), '.env')}")
        else:
            print(f"[SUCCESS] GROQ_API_KEY loaded (length: {len(groq_api_key)}, ends with: ...{groq_api_key[-4:]})")
            if len(groq_api_key) < 50:
                print("[WARNING] WARNING: GROQ API key seems too short. Please check if it's complete.")
            
        try:
            self.llm = ChatGroq(
                model="gemma2-9b-it",
                api_key=groq_api_key,
                temperature=0.3
            )
            print("[SUCCESS] GROQ LLM client initialized")
        except Exception as e:
            print(f"[ERROR] Failed to initialize GROQ LLM: {str(e)}")
            self.llm = None
        
        self.vector_store: Optional[FAISS] = None
        
        if self.embeddings_available and os.path.exists(self.faiss_index_path):
            try:
                self.vector_store = FAISS.load_local(
                    self.faiss_index_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                print("[DOCS] Loading Math and English modules...")
            except Exception as e:
                print(f"[WARNING] Failed to load existing index: {str(e)}")
                self.vector_store = None
        else:
            if self.embeddings_available:
                print("[DOCS] Ready for document ingestion...")
            else:
                print("[DOCS] Running in offline mode - limited functionality available")

    def _load_embeddings_with_retry(self, max_retries=3) -> tuple:
        """Load embeddings with retry mechanism and cache clearing."""
        model_name = "sentence-transformers/all-MiniLM-L6-v2"
        
        for attempt in range(max_retries):
            try:
                print(f"[INFO] Loading embeddings model (attempt {attempt + 1}/{max_retries})...")
                
                # Clear cache if this isn't the first attempt
                if attempt > 0:
                    self._clear_model_cache(model_name)
                    time.sleep(2)  # Brief delay before retry
                
                embeddings = HuggingFaceEmbeddings(
                    model_name=model_name,
                    model_kwargs={"device": "cpu"},
                    encode_kwargs={"normalize_embeddings": True}
                )
                
                print("[SUCCESS] Embeddings model loaded successfully")
                return embeddings, True
                
            except Exception as e:
                print(f"[WARNING] Attempt {attempt + 1} failed: {str(e)}")
                if attempt < max_retries - 1:
                    print(f"[INFO] Retrying in 3 seconds...")
                    time.sleep(3)
                else:
                    print("[ERROR] All attempts failed. Running in offline mode.")
        
        return None, False
    
    def _clear_model_cache(self, model_name: str):
        """Clear HuggingFace model cache to force fresh download."""
        try:
            # Try to find and clear HF cache
            cache_dirs = [
                os.path.expanduser("~/.cache/huggingface/transformers"),
                os.path.expanduser("~/.cache/huggingface/hub"),
            ]
            
            for cache_dir in cache_dirs:
                if os.path.exists(cache_dir):
                    # Look for model-specific cached files
                    for item in os.listdir(cache_dir):
                        if "all-MiniLM-L6-v2" in item or "sentence-transformers" in item:
                            item_path = os.path.join(cache_dir, item)
                            if os.path.isdir(item_path):
                                print(f"[CLEAR] Clearing cache: {item_path}")
                                shutil.rmtree(item_path, ignore_errors=True)
                            elif os.path.isfile(item_path):
                                os.remove(item_path)
        except Exception as e:
            print(f"[WARNING] Could not clear cache: {e}")
    
    def retry_embeddings_loading(self) -> bool:
        """Manual retry method for embeddings loading."""
        print("[INFO] Manually retrying embeddings loading...")
        self.embeddings, self.embeddings_available = self._load_embeddings_with_retry()
        
        # Try to reload vector store if embeddings are now available
        if self.embeddings_available and os.path.exists(self.faiss_index_path):
            try:
                self.vector_store = FAISS.load_local(
                    self.faiss_index_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                print("[SUCCESS] Vector store reloaded successfully")
            except Exception as e:
                print(f"[WARNING] Failed to reload vector store: {str(e)}")
        
        return self.embeddings_available
    
    def get_document_count(self) -> int:
        return self.vector_store.index.ntotal if self.vector_store else 0

    def detect_subject_and_intent(self, question: str) -> Dict[str, str]:
        question_lower = question.lower().strip()
        
        subject = "general"
        if any(word in question_lower for word in ["math", "number", "count", "add", "subtract", "plus", "minus", "multiply", "divide"]):
            subject = "math"
        elif any(word in question_lower for word in ["read", "story", "letter", "word", "english", "spell"]):
            subject = "reading"
        
        intent = "question"
        if any(phrase in question_lower for phrase in ["teach me", "learn", "help me", "show me how"]):
            intent = "learn"
        elif any(phrase in question_lower for phrase in ["what's in", "chapter", "lesson", "unit"]):
            intent = "explore"
        elif any(phrase in question_lower for phrase in ["practice", "exercise", "quiz", "test"]):
            intent = "practice"
        
        return {"subject": subject, "intent": intent}

    def should_use_rag(self, question: str) -> bool:
        if not self.embeddings_available:
            return False
            
        analysis = self.detect_subject_and_intent(question)
        
        if self.vector_store and self.get_document_count() > 0:
            if analysis["intent"] in ["learn", "explore", "practice"]:
                return True
            if analysis["subject"] in ["math", "reading"]:
                return True
        
        greetings = ["hello", "hi", "thanks", "bye", "good morning", "how are you"]
        if any(greeting in question.lower() for greeting in greetings):
            return False
        
        if re.search(r'\d+\s*[+\-*/]\s*\d+', question):
            return False
        
        return self.vector_store is not None and self.get_document_count() > 0

    def evaluate_simple_math(self, question: str) -> Optional[str]:
        math_pattern = r'(\d+)\s*([+\-*/])\s*(\d+)'
        match = re.search(math_pattern, question)
        
        if match:
            num1, operator, num2 = match.groups()
            num1, num2 = int(num1), int(num2)
            
            operations = {
                '+': num1 + num2,
                '-': num1 - num2,
                '*': num1 * num2,
                '/': num1 / num2 if num2 != 0 else "Cannot divide by zero!"
            }
            
            result = operations.get(operator)
            if isinstance(result, (int, float)):
                return f"Great math question! Let me help you solve {num1} {operator} {num2}:\n\n{num1} {operator} {num2} = {result}\n\nExcellent work! Would you like to try another math problem? 🎉"
            else:
                return str(result)
        return None

    def get_conversation_context(self) -> str:
        if len(self.conversation_history) <= 1:
            return ""
        
        recent_history = self.conversation_history[-20:]
        history_text = "\n".join([
            f"{'Student' if msg['role'] == 'user' else 'Teacher'}: {msg['content']}"
            for msg in recent_history[:-1]
        ])
        
        context = f"\nRecent conversation:\n{history_text}\n"
        
        if self.current_subject:
            context += f"\nCurrent subject focus: {self.current_subject}\n"
        
        return context

    def get_relevant_context(self, question: str, subject: str, top_k: int = 5) -> List[Document]:
        if not self.vector_store:
            return []
        
        all_docs = self.vector_store.similarity_search(question, k=top_k*2)
        
        if subject != "general":
            filtered_docs = [doc for doc in all_docs if doc.metadata.get("subject") == subject]
            if filtered_docs:
                return filtered_docs[:top_k]
        
        return all_docs[:top_k]

    def create_educational_prompt(self, question: str, context_docs: List[Document], analysis: Dict[str, str]) -> str:
        conversation_context = self.get_conversation_context()
        
        if context_docs:
            docs_formatted = "\n\n".join([
                f"[DOCS] Source: {doc.metadata.get('source', 'Unknown')}\n{doc.page_content}"
                for doc in context_docs
            ])
            
            if analysis["intent"] == "learn":
                prompt = f"""You are an enthusiastic and patient teacher for 1st and 2nd grade students. 

{conversation_context}

EDUCATIONAL CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Please provide a structured lesson response:
1. Start with encouragement
2. Break down the concept into simple steps
3. Use examples from the educational content
4. Include interactive elements (questions, activities)
5. End with positive reinforcement and next steps

Keep the language simple and engaging for young learners. 
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems."""

            elif analysis["intent"] == "explore":
                prompt = f"""You are a friendly teacher helping students explore educational content.

{conversation_context}

AVAILABLE CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Give an overview that:
1. Summarizes what's available in an exciting way
2. Highlights key learning objectives
3. Suggests where to start
4. Makes it sound fun and achievable

Keep the language simple and engaging for young learners. 
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.
"""

            else:
                prompt = f"""You are a helpful teacher for young students.

{conversation_context}

EDUCATIONAL CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Provide a clear, encouraging answer using the educational content. Keep it simple and age-appropriate for 1st/2nd graders.
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.
"""

        else:
            prompt = f"""You are a warm, encouraging teacher for 1st and 2nd grade students.

{conversation_context}

STUDENT QUESTION: {question}

Provide a helpful, age-appropriate response using your general knowledge. Be encouraging and suggest how they might learn more about this topic.
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.
"""

        return prompt

    def ingest_file(self, file_name: str, file_bytes: bytes) -> str:
        if not self.embeddings_available:
            # Try to reload embeddings once
            print("[INFO] Embeddings not available, attempting to reload...")
            if not self.retry_embeddings_loading():
                return f"❗ Cannot ingest '{file_name}': Embeddings model not available. Check internet connection."
        
        try:
            suffix = os.path.splitext(file_name)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            
            new_raw_docs = load_document_from_path(tmp_path)
            os.remove(tmp_path)
            
            if not new_raw_docs:
                return f"'{file_name}' uploaded, but no content found."
            
            new_chunks = smart_chunk_documents(new_raw_docs)
            
            if not new_chunks:
                return f"'{file_name}' uploaded, but no chunks created."
            
            if self.vector_store is None:
                self.vector_store = FAISS.from_documents(new_chunks, self.embeddings)
            else:
                self.vector_store.add_documents(new_chunks)
            
            self.vector_store.save_local(self.faiss_index_path)
            
            destination = os.path.join(self.doc_folder, file_name)
            with open(destination, "wb") as f:
                f.write(file_bytes)
            
            return f"📥 '{file_name}' ingested. Added {len(new_chunks)} new chunk(s)."
        
        except Exception as e:
            return f"❗ Failed to ingest '{file_name}': {str(e)}"

    def clear_all_data(self) -> str:
        self.vector_store = None
        self.conversation_history = []
        self.current_subject = None
        self.learning_progress = {}
        
        for folder in [self.doc_folder, self.index_folder]:
            for path in glob.glob(os.path.join(folder, "*")):
                try:
                    if os.path.isfile(path) or os.path.islink(path):
                        os.remove(path)
                    elif os.path.isdir(path):
                        import shutil
                        shutil.rmtree(path)
                except Exception as e:
                    pass
        
        return "[CLEAR] Cache cleared. All documents, indexes, and conversation history removed."

    def query(self, question: str, top_k: int = 5) -> str:
        if not question:
            return "Please type a question."
        
        question_clean = question.strip().lower()
        if question_clean in {"what", "why", "how", "where", "when", "ok", "yes", "no"}:
            return "Could you tell me a bit more about what you want to know? I'm here to help you learn!"
        
        self.conversation_history.append({"role": "user", "content": question})
        
        math_result = self.evaluate_simple_math(question)
        if math_result:
            self.conversation_history.append({"role": "assistant", "content": math_result})
            return math_result
        
        analysis = self.detect_subject_and_intent(question)
        
        if analysis["subject"] == "general" and self.current_subject:
            analysis["subject"] = self.current_subject
        else:
            self.current_subject = analysis["subject"]
        
        use_rag = self.should_use_rag(question)
        
        if use_rag:
            context_docs = self.get_relevant_context(question, analysis["subject"], top_k)
            prompt = self.create_educational_prompt(question, context_docs, analysis)
        else:
            conversation_context = self.get_conversation_context()
            prompt = f"""You are a friendly, patient teacher for 1st and 2nd grade students.

{conversation_context}

STUDENT QUESTION: {question}

Provide an encouraging, age-appropriate response. Use simple language and make learning fun! 
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.
"""
        
        if not self.llm:
            answer = "[ERROR] Language model is not available. Please check your GROQ_API_KEY in the .env file and restart the application."
        else:
            try:
                response = self.llm.invoke(prompt)
                answer = response.content.strip()
            except Exception as e:
                print(f"[WARNING] LLM Error: {str(e)}")
                if "401" in str(e) or "invalid" in str(e).lower():
                    answer = "[ERROR] Invalid GROQ API key. Please check your GROQ_API_KEY in backend/.env file. Get a new key from https://console.groq.com/keys"
                elif not self.embeddings_available:
                    # Try to reload embeddings once per session
                    if not hasattr(self, '_embeddings_retry_attempted'):
                        self._embeddings_retry_attempted = True
                        print("[INFO] Attempting to reload embeddings...")
                        if self.retry_embeddings_loading():
                            # Recursive call with embeddings now available
                            return self.query(question, top_k)
                    
                    answer = "I'm running in offline mode right now. I can help with simple math problems like '5 + 3' or general conversations, but I can't access educational documents. Check your internet connection and try restarting the application."
                else:
                    answer = f"I'm having trouble connecting to my language model right now: {str(e)}. Please try again or check your API key."
        
        self.conversation_history.append({"role": "assistant", "content": answer})
        
        if len(self.conversation_history) > 24:
            self.conversation_history = self.conversation_history[-24:]
        
        return answer
