import os
import glob
import tempfile
import re
import time
import shutil
from typing import List, Optional, Dict, Any
from dotenv import load_dotenv

# Document Processing Imports
import pypdf
import docx2txt
from pptx import Presentation

# LangChain / AI Imports
from langchain_groq import ChatGroq
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# Load environment variables
load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))


def extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a PDF file."""
    try:
        reader = pypdf.PdfReader(file_path)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception as e:
        print(f"[ERROR] PDF Extraction failed: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    try:
        return docx2txt.process(file_path)
    except Exception as e:
        print(f"[ERROR] DOCX Extraction failed: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)
    except Exception as e:
        print(f"[ERROR] PPTX Extraction failed: {e}")
        return ""


def extract_text_from_txt(file_path: str) -> str:
    """Extracts text from a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"[ERROR] TXT Extraction failed: {e}")
        return ""


def load_document_from_path(file_path: str) -> List[Dict[str, Any]]:
    """Determines file type and extracts text accordingly."""
    ext = os.path.splitext(file_path)[1].lower()
    docs = []
    filename = os.path.basename(file_path)
    content = ""

    try:
        if ext == ".pdf":
            content = extract_text_from_pdf(file_path)
        elif ext == ".docx":
            content = extract_text_from_docx(file_path)
        elif ext == ".pptx":
            content = extract_text_from_pptx(file_path)
        elif ext == ".txt":
            content = extract_text_from_txt(file_path)
        else:
            return docs

        if content.strip():
            docs.append({
                "content": content,
                "filename": filename,
                # Basic classification based on filename
                "subject": "math" if "aejm" in filename.lower()
                else "reading" if "aemr" in filename.lower()
                else "general"
            })
    except Exception as e:
        print(f"[ERROR] Loading document failed: {e}")

    return docs


def smart_chunk_documents(documents: List[Dict[str, Any]]) -> List[Document]:
    """Chunks documents based on logical sections (Chapters, Lessons) and size."""
    chunks = []

    for doc_data in documents:
        content = doc_data["content"]
        filename = doc_data["filename"]
        subject = doc_data["subject"]

        # Split by logical headers using lookahead regex
        sections = re.split(r'\n(?=Chapter|Lesson|Unit|Exercise|Activity|\d+\.)', content)

        for section_idx, section in enumerate(sections):
            # Skip very short snippets
            if len(section.strip()) < 50:
                continue

            # If section is too long, sub-chunk it
            if len(section) > 1000:
                words = section.split()
                # Overlapping windows
                for i in range(0, len(words), 300):
                    chunk_words = words[i:i + 400]
                    chunk_text = " ".join(chunk_words)

                    if len(chunk_text.strip()) > 30:
                        doc_obj = Document(
                            page_content=chunk_text,
                            metadata={
                                "source": filename,
                                "subject": subject,
                                "section": section_idx,
                                "chunk_type": "content"
                            }
                        )
                        chunks.append(doc_obj)
            else:
                doc_obj = Document(
                    page_content=section,
                    metadata={
                        "source": filename,
                        "subject": subject,
                        "section": section_idx,
                        "chunk_type": "section"
                    }
                )
                chunks.append(doc_obj)

    return chunks


class RAGSystem:
    def __init__(self, doc_folder: str = "./docs", index_folder: str = "./indexes"):
        self.doc_folder = doc_folder
        self.index_folder = index_folder
        self.faiss_index_path = os.path.join(self.index_folder, "faiss_index")

        self.conversation_history = []
        self.current_subject = None
        self.learning_progress = {}

        os.makedirs(self.doc_folder, exist_ok=True)
        os.makedirs(self.index_folder, exist_ok=True)

        self.embeddings, self.embeddings_available = self._load_embeddings_with_retry()

        groq_api_key = os.getenv("GROQ_API_KEY")
        if not groq_api_key:
            print("[WARNING] GROQ_API_KEY not found in environment variables")
            print(f"Looking for .env file at: {os.path.join(os.path.dirname(__file__), '.env')}")
        else:
            safe_key_preview = f"...{groq_api_key[-4:]}" if len(groq_api_key) > 4 else "Invalid"
            print(f"[SUCCESS] GROQ_API_KEY loaded (length: {len(groq_api_key)}, ends with: {safe_key_preview})")
            if len(groq_api_key) < 50:
                print("[WARNING] GROQ API key seems too short. Please check if it's complete.")

        try:
            self.llm = ChatGroq(
                model="openai/gpt-oss-120b",
                api_key=groq_api_key,
                temperature=0.3
            )
            print("[SUCCESS] GROQ LLM client initialized")
        except Exception as e:
            print(f"[ERROR] Failed to initialize GROQ LLM: {str(e)}")
            self.llm = None

        self.vector_store: Optional[FAISS] = None

        if self.embeddings_available and os.path.exists(self.faiss_index_path):
            try:
                self.vector_store = FAISS.load_local(
                    self.faiss_index_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                print(f"[DOCS] Loaded existing index with {self.vector_store.index.ntotal} chunks.")
            except Exception as e:
                print(f"[WARNING] Failed to load existing index: {e}")
                self.vector_store = None
        else:
            if self.embeddings_available:
                print("[DOCS] Ready for document ingestion...")
            else:
                print("[DOCS] Running in offline mode - limited functionality available")

    def _language_display_name(self, language: Optional[str]) -> str:
        mapping = {"en": "English", "ta": "Tamil"}
        if not language:
            return "English"
        return mapping.get(language.lower(), "English")

    def _contains_tamil(self, text: str) -> bool:
        """Return True if the string contains any Tamil Unicode characters."""
        for ch in text:
            code = ord(ch)
            if 0x0B80 <= code <= 0x0BFF:
                return True
        return False

    def _build_language_instruction(self, target_language: str) -> str:
        """
        Return language and style rules depending on target language.

        For Tamil, this includes strong guardrails to avoid phonics /
        worksheet-style responses unless explicitly requested, and to
        avoid irrelevant 'tree sharing fruits' reading passages.
        """
        lang = (target_language or "en").lower()
        language_name = self._language_display_name(lang)

        if lang == "ta":
            return (
                "LANGUAGE AND STYLE RULES\n"
                "- Answer only in Tamil.\n"
                "- Use simple, short Tamil sentences that a Grade 1 or Grade 2 child can understand.\n"
                "- Use everyday Tamil words. Avoid very complex or literary Tamil.\n"
                "- Do not mix English words unless they are proper names or absolutely necessary for the lesson content.\n"
                "- Do not create alphabet or phonics drills such as 'A is for ant' or 'B is for bag' unless the student clearly asks you to teach letters or phonics.\n"
                "- Do not create classroom worksheet instructions such as asking the student to circle words, repeat after you, sit in a circle, or similar activities unless the student clearly asks for such activities.\n"
                "- When the question asks about an animal, place, object, person, story, or a topic like 'கூட்டல் பற்றி', give a direct, factual explanation in simple Tamil instead of turning it into an exercise.\n"
                "- The student has ALREADY asked a question. Do NOT answer by telling them to ask a question again.\n"
                "- Specifically, do NOT use sentences like:\n"
                "  'உனக்கு என்ன தெரியவில்லை',\n"
                "  'நீ என்ன கற்றுக்கொள்ள விரும்புகிறாய்',\n"
                "  'நீ என்ன பற்றி கேட்க விரும்புகிறாய்',\n"
                "  'உன் கேள்வியை எழுது',\n"
                "  'கேள்வியை கேள்',\n"
                "  'கேள்வியை கேட்டால் மட்டுமே நான் பதில் அளிக்க முடியும்'.\n"
                "- If the question is about math (for example it mentions 'கூட்டல்', 'கழித்தல்', 'பெருக்கல்', or 'வகுத்தல்'), you MUST explain the math idea directly.\n"
                "- For math questions, do NOT talk about trees sharing their fruits, children sharing fruits, poems about 'for' and 'on', or similar reading-passage content unless the student explicitly asks about trees, fruits, or that poem.\n"
                "- Your FIRST sentence must start explaining the topic in the student's question. It must not be a question back to the student.\n"
                "- Keep the answer as clean plain text with no emojis and no visible formatting marks, so it is easy to use with text to speech."
            )
        else:
            return (
                "LANGUAGE AND STYLE RULES\n"
                f"- Answer only in {language_name}.\n"
                "- Use simple, short sentences that a 1st or 2nd grade child can understand.\n"
                "- Explain ideas clearly and directly. Avoid long, complex sentences.\n"
                "- First, answer the student's question as clearly as you can. Only after answering may you add one short suggestion or follow-up question.\n"
                "- Do not reply with vague questions like 'What do you want to learn?' unless the student clearly asks for help choosing a topic.\n"
                "- Do not turn everything into a quiz or worksheet unless the student asks for practice.\n"
                "- Keep the answer as clean plain text with no emojis and no visible formatting marks, so it is easy to use with text to speech."
            )

    def _load_embeddings_with_retry(self, max_retries=3) -> tuple:
        """Load embeddings with retry mechanism and cache clearing."""
        model_name = "sentence-transformers/all-MiniLM-L6-v2"

        for attempt in range(max_retries):
            try:
                print(f"[INFO] Loading embeddings model (attempt {attempt + 1}/{max_retries})...")

                # Clear cache if this isn't the first attempt
                if attempt > 0:
                    self._clear_model_cache(model_name)
                    time.sleep(2)  # Brief delay before retry

                embeddings = HuggingFaceEmbeddings(
                    model_name=model_name,
                    model_kwargs={"device": "cpu"},
                    encode_kwargs={"normalize_embeddings": True}
                )

                print("[SUCCESS] Embeddings model loaded successfully")
                return embeddings, True

            except Exception as e:
                print(f"[WARNING] Attempt {attempt + 1} failed: {e}")
                if attempt < max_retries - 1:
                    print(f"[INFO] Retrying in 3 seconds...")
                    time.sleep(3)
                else:
                    print("[ERROR] All attempts failed. Running in offline mode.")

        return None, False

    def _clear_model_cache(self, model_name: str):
        """Clear HuggingFace model cache to force fresh download."""
        try:
            cache_dirs = [
                os.path.expanduser("~/.cache/huggingface/transformers"),
                os.path.expanduser("~/.cache/huggingface/hub"),
            ]

            for cache_dir in cache_dirs:
                if os.path.exists(cache_dir):
                    for item in os.listdir(cache_dir):
                        if "all-MiniLM-L6-v2" in item or "sentence-transformers" in item:
                            item_path = os.path.join(cache_dir, item)
                            if os.path.isdir(item_path):
                                print(f"[CLEAR] Clearing cache: {item_path}")
                                shutil.rmtree(item_path, ignore_errors=True)
                            elif os.path.isfile(item_path):
                                os.remove(item_path)
        except Exception as e:
            print(f"[WARNING] Could not clear cache: {e}")

    def retry_embeddings_loading(self) -> bool:
        """Manual retry method for embeddings loading."""
        print("[INFO] Manually retrying embeddings loading...]")
        self.embeddings, self.embeddings_available = self._load_embeddings_with_retry()

        # Try to reload vector store if embeddings are now available
        if self.embeddings_available and os.path.exists(self.faiss_index_path):
            try:
                self.vector_store = FAISS.load_local(
                    self.faiss_index_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                print("[SUCCESS] Vector store reloaded successfully")
            except Exception as e:
                print(f"[WARNING] Failed to reload vector store: {e}")

        return self.embeddings_available

    def get_document_count(self) -> int:
        return self.vector_store.index.ntotal if self.vector_store else 0

    def detect_subject_and_intent(self, question: str) -> Dict[str, str]:
        """
        Heuristic subject/intent detection for English + Tamil.
        """
        question_lower = question.lower().strip()
        q = question.strip()

        subject = "general"

        # English subject hints
        if any(word in question_lower for word in ["math", "number", "count", "add", "subtract", "plus", "minus", "multiply", "divide"]):
            subject = "math"
        elif any(word in question_lower for word in ["read", "story", "letter", "word", "english", "spell"]):
            subject = "reading"

        # Tamil subject hints
        if any(tok in q for tok in ["கூட்டல்", "கழித்தல்", "பெருக்கல்", "வகுத்தல்"]):
            subject = "math"
        if any(tok in q for tok in ["கதை", "கவிதை", "கவிதையில்", "வாசிப்பு"]):
            if subject == "general":
                subject = "reading"

        intent = "question"

        # English intent
        if any(phrase in question_lower for phrase in ["teach me", "learn", "help me", "show me how"]):
            intent = "learn"
        elif any(phrase in question_lower for phrase in ["what's in", "chapter", "lesson", "unit"]):
            intent = "explore"
        elif any(phrase in question_lower for phrase in ["practice", "exercise", "quiz", "test"]):
            intent = "practice"

        # Tamil intent
        if any(tok in q for tok in ["கற்றுக்கொடு", "கற்றுக்கொடுங்கள்", "கற்று கொடு", "விளக்கவும்"]):
            intent = "learn"
        if any(tok in q for tok in ["பயிற்சி", "வினா", "வினாத்தாள்"]):
            if intent == "question":
                intent = "practice"

        return {"subject": subject, "intent": intent}

    def should_use_rag(self, question: str) -> bool:
        """
        Decide whether to use RAG.

        IMPORTANT: For Tamil questions, we ALWAYS skip RAG and rely
        only on the model + prompts. This avoids pulling in unrelated
        Tamil textbook passages.
        """
        if not self.embeddings_available or not self.vector_store:
            return False

        # If the question contains Tamil script, do not use RAG.
        if self._contains_tamil(question):
            return False

        analysis = self.detect_subject_and_intent(question)

        if self.vector_store and self.get_document_count() > 0:
            if analysis["intent"] in ["learn", "explore", "practice"]:
                return True
            if analysis["subject"] in ["math", "reading"] and self.vector_store:
                return True

        greetings = ["hello", "hi", "thanks", "bye", "good morning", "how are you"]
        if any(greeting in question.lower() for greeting in greetings):
            return False

        if re.search(r'\d+\s*[+\-*/]\s*\d+', question):
            return False

        return self.vector_store is not None and self.get_document_count() > 0

    def evaluate_simple_math(self, question: str, target_language: str = "en") -> Optional[str]:
        """
        Handle simple arithmetic, including:
        - Digit forms: 2+2, 2 + 2, 2×2, 2÷2, etc.
        - Basic Tamil word forms: 'இரண்டு கூட்டி இரண்டு', 'இரண்டு ப்ளஸ் இரண்டு', etc.
        """
        # 1) Try numeric pattern first (handles "2+2", "2 + 2", maybe "2×2", "2÷2")
        math_pattern = r'(\d+)\s*([+\-x×*/÷])\s*(\d+)'
        match = re.search(math_pattern, question)

        def _compute(a: int, op_symbol: str, b: int):
            if op_symbol in ["+", "＋"]:
                return a + b
            if op_symbol in ["-", "−"]:
                return a - b
            if op_symbol in ["*", "x", "×"]:
                return a * b
            if op_symbol in ["/", "÷"]:
                if b == 0:
                    return None
                return a / b
            return None

        def _format_result(a: int, op_symbol: str, b: int, result_value):
            if isinstance(result_value, float) and result_value.is_integer():
                result_value = int(result_value)

            if target_language == "ta":
                if op_symbol in ["+", "＋"]:
                    op_word = "கூட்டல்"
                elif op_symbol in ["-", "−"]:
                    op_word = "கழித்தல்"
                elif op_symbol in ["*", "x", "×"]:
                    op_word = "பெருக்கல்"
                elif op_symbol in ["/", "÷"]:
                    op_word = "வகுத்தல்"
                else:
                    op_word = "கணக்கு"
                return f"{a} {op_symbol} {b} = {result_value}. இது ஒரு எளிய {op_word} எடுத்துக்காட்டு."
            else:
                return f"{a} {op_symbol} {b} = {result_value}. This is a simple example of arithmetic."

        if match:
            num1_str, operator, num2_str = match.groups()
            num1, num2 = int(num1_str), int(num2_str)
            result = _compute(num1, operator, num2)
            if result is None:
                if target_language == "ta":
                    return "பூஜ்யம் மூலம் பகுக்க முடியாது."
                else:
                    return "I cannot divide by zero!"
            return _format_result(num1, operator, num2, result)

        # 2) If no digit-based match, try Tamil-word based patterns.
        # Map a few common Tamil number words (good enough for Grade 1–2).
        TAMIL_NUM_WORDS = {
            "பூஜ்யம்": 0,
            "சூன்யம்": 0,
            "ஒன்று": 1,
            "ஒரு": 1,
            "இரண்டு": 2,
            "மூன்று": 3,
            "நான்கு": 4,
            "ஐந்து": 5,
            "ஆறு": 6,
            "ஏழு": 7,
            "எட்டு": 8,
            "ஒன்பது": 9,
            "பத்து": 10,
        }

        def tamil_to_int(word: str) -> Optional[int]:
            w = word.strip()
            return TAMIL_NUM_WORDS.get(w)

        # Common Tamil operator words
        OP_WORDS = {
            "+": ["கூட்ட", "கூட்டி", "கூட்டினால்", "ப்ளஸ்", "பிளஸ்"],
            "-": ["கழித்த", "கழித்து", "குறை"],
            "*": ["பெருக்கு", "பெருக்கி", "மடங்கு"],
            "/": ["வகுத்து", "வகுத்தல்", "பகுத்து"],
        }

        # Normalize question a bit (remove question tail words like "என்ன", "எவ்வளவு", "என்றால்")
        q_norm = question.replace("?", " ")
        for tail in ["என்ன", "எவ்வளவு", "எவ்வளவு?", "என்றால்"]:
            q_norm = q_norm.replace(tail, " ")
        tokens = [t for t in re.split(r"\s+", q_norm) if t]

        # Try pattern: [num-word] [op-word] [num-word]
        for i in range(len(tokens) - 2):
            w1, wop, w2 = tokens[i], tokens[i + 1], tokens[i + 2]
            n1 = tamil_to_int(w1)
            n2 = tamil_to_int(w2)
            if n1 is None or n2 is None:
                continue

            op_symbol = None
            for sym, word_list in OP_WORDS.items():
                if any(wop.startswith(ow) for ow in word_list):
                    op_symbol = sym
                    break

            if op_symbol:
                result = _compute(n1, op_symbol, n2)
                if result is None:
                    if target_language == "ta":
                        return "பூஜ்யம் மூலம் பகுக்க முடியாது."
                    else:
                        return "I cannot divide by zero!"
                return _format_result(n1, op_symbol, n2, result)

        return None

    def get_conversation_context(self) -> str:
        if len(self.conversation_history) <= 1:
            return ""

        recent_history = self.conversation_history[-20:]
        history_text = "\n".join([
            f"{'Student' if msg['role'] == 'user' else 'Teacher'}: {msg['content']}"
            for msg in recent_history[:-1]
        ])

        context = f"\nRecent conversation:\n{history_text}\n"

        if self.current_subject:
            context += f"\nCurrent subject focus: {self.current_subject}\n"

        return context

    def get_relevant_context(self, question: str, subject: str, top_k: int = 5) -> List[Document]:
        if not self.vector_store:
            return []

        try:
            all_docs = self.vector_store.similarity_search(question, k=top_k * 2)

            if subject != "general":
                filtered_docs = [doc for doc in all_docs if doc.metadata.get("subject") == subject]
                if filtered_docs:
                    return filtered_docs[:top_k]

            return all_docs[:top_k]
        except Exception as e:
            print(f"[WARNING] Similarity search failed: {e}")
            return []

    def create_educational_prompt(
        self,
        question: str,
        context_docs: List[Document],
        analysis: Dict[str, str],
        target_language: str = "en",
    ) -> str:
        conversation_context = self.get_conversation_context()
        language_instruction = self._build_language_instruction(target_language)

        docs_formatted = ""
        if context_docs:
            docs_formatted = "\n\n".join([
                f"[DOCS] Source: {doc.metadata.get('source', 'Unknown')}\n{doc.page_content}"
                for doc in context_docs
            ])

            if analysis["intent"] == "learn":
                prompt = f"""You are an enthusiastic and patient teacher for 1st and 2nd grade students.

GENERAL BEHAVIOUR
- Always answer clearly, using simple words and short sentences.
- Give a direct answer to the student's question first.
- Use the educational content only when it is relevant and helpful.
- Do not start by asking the student what they want to learn; they already asked a question.

{conversation_context}

EDUCATIONAL CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Please provide a structured lesson response:
1. Start with brief encouragement in the student's language.
2. Directly answer the student's question in a simple way.
3. Break down the concept into simple steps.
4. Use examples from the educational content when they are helpful.
5. You may ask one or two simple follow-up questions, but do not turn the entire answer into a written worksheet.
6. End with gentle positive reinforcement and a simple suggestion for what they could learn next.

Keep the language simple and engaging for young learners.
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.

{language_instruction}
"""

            elif analysis["intent"] == "explore":
                prompt = f"""You are a friendly teacher helping students explore educational content.

GENERAL BEHAVIOUR
- Always answer clearly, using simple words and short sentences.
- Give an overview that is encouraging and not overwhelming.
- Use the educational content only when it is relevant.
- Do not reply with vague questions; always give some concrete information first.

{conversation_context}

AVAILABLE CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Give an overview that:
1. Summarizes what is available in a simple, exciting way.
2. Highlights a few key learning objectives in plain language.
3. Suggests a clear and simple place to start.
4. Encourages the student and makes the topic feel achievable.

Keep the language simple and engaging for young learners.
Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.

{language_instruction}
"""

            else:
                prompt = f"""You are a helpful teacher for young students.

GENERAL BEHAVIOUR
- Always answer clearly, using simple words and short sentences.
- Focus on directly answering the student's question first.
- Use the educational content only when it is helpful.
- Do not turn the answer into a worksheet unless the student clearly asked for practice.

{conversation_context}

EDUCATIONAL CONTENT:
{docs_formatted}

STUDENT QUESTION: {question}

Provide a clear, encouraging answer using the educational content.
Keep it simple and age-appropriate for 1st and 2nd graders.
Answer the question directly before adding any extra suggestions.

Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.

{language_instruction}
"""

        else:
            prompt = f"""You are a warm, encouraging teacher for 1st and 2nd grade students.

GENERAL BEHAVIOUR
- Always answer clearly, using simple words and short sentences.
- Give a direct, helpful answer to the student's question first.
- Encourage curiosity and make the student feel confident.
- Do not reply with vague questions like "What do you want to learn?" because the student has already asked a question.

{conversation_context}

STUDENT QUESTION: {question}

Provide a helpful, age-appropriate response using your general knowledge.
First, answer the student's question as clearly as you can.
Then, if helpful, suggest one simple way they might learn more about this topic.

Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.

{language_instruction}
"""

        return prompt

    def ingest_file(self, file_name: str, file_bytes: bytes) -> str:
        if not self.embeddings_available:
            # Try to reload embeddings once
            print("[INFO] Embeddings not available, attempting to reload...]")
            if not self.retry_embeddings_loading():
                return f"❗ Cannot ingest '{file_name}': Embeddings model not available. Check internet connection."

        try:
            suffix = os.path.splitext(file_name)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            new_raw_docs = load_document_from_path(tmp_path)

            # Clean up temp file
            try:
                os.remove(tmp_path)
            except Exception:
                pass

            if not new_raw_docs:
                return f"'{file_name}' uploaded, but no content found."

            new_chunks = smart_chunk_documents(new_raw_docs)

            if not new_chunks:
                return f"'{file_name}' uploaded, but no chunks created."

            if self.vector_store is None:
                self.vector_store = FAISS.from_documents(new_chunks, self.embeddings)
            else:
                self.vector_store.add_documents(new_chunks)

            # Ensure index folder exists before saving
            os.makedirs(self.index_folder, exist_ok=True)
            self.vector_store.save_local(self.faiss_index_path)

            # Ensure docs folder exists
            os.makedirs(self.doc_folder, exist_ok=True)
            destination = os.path.join(self.doc_folder, file_name)
            with open(destination, "wb") as f:
                f.write(file_bytes)

            return f"📥 '{file_name}' ingested. Added {len(new_chunks)} new chunk(s)."

        except Exception as e:
            return f"❗ Failed to ingest '{file_name}': {e}"

    def clear_all_data(self) -> str:
        self.vector_store = None
        self.conversation_history = []
        self.current_subject = None
        self.learning_progress = {}

        for folder in [self.doc_folder, self.index_folder]:
            if not os.path.exists(folder):
                continue
            for path in glob.glob(os.path.join(folder, "*")):
                try:
                    if os.path.isfile(path) or os.path.islink(path):
                        os.remove(path)
                    elif os.path.isdir(path):
                        shutil.rmtree(path)
                except Exception as e:
                    print(f"[ERROR] Failed to delete {path}: {e}")

        return "[CLEAR] Cache cleared. All documents, indexes, and conversation history removed."

    def query(self, question: str, top_k: int = 5, target_language: str = "en") -> str:
        if not question:
            return "Please type a question."

        question_clean = question.strip().lower()
        if question_clean in {"what", "why", "how", "where", "when", "ok", "yes", "no"}:
            return "Could you tell me a bit more about what you want to know? I'm here to help you learn!"

        self.conversation_history.append({"role": "user", "content": question})

        analysis = self.detect_subject_and_intent(question)

        if analysis["subject"] == "general" and self.current_subject:
            analysis["subject"] = self.current_subject
        else:
            self.current_subject = analysis["subject"]

        normalized_language = (target_language or "en").lower()
        if normalized_language not in {"en", "ta"}:
            normalized_language = "en"

        # Try to handle simple math (both digit-based and Tamil word-based) BEFORE calling the LLM
        math_result = self.evaluate_simple_math(question, target_language=normalized_language)
        if math_result:
            self.conversation_history.append({"role": "assistant", "content": math_result})
            if len(self.conversation_history) > 24:
                self.conversation_history = self.conversation_history[-24:]
            return math_result

        use_rag = self.should_use_rag(question)
        is_tamil = self._contains_tamil(question)

        if use_rag:
            context_docs = self.get_relevant_context(question, analysis["subject"], top_k)
            prompt = self.create_educational_prompt(
                question,
                context_docs,
                analysis,
                target_language=normalized_language,
            )
        else:
            language_instruction = self._build_language_instruction(normalized_language)

            if is_tamil:
                # Tamil factual explainer path – NO alphabet songs / worksheets / tree-fruit poems
                prompt = f"""You are a factual explanation assistant for young children in Grades 1 and 2.

ROLE
- Your job is to explain clearly what the student asked about.
- You are NOT running a classroom. You are NOT giving phonics drills or alphabet lessons unless the question is about letters or reading.
- You are NOT asking the child what they want to learn. You must treat the given sentence as their question.

LANGUAGE RULES
- Answer only in Tamil.
- Use simple, short Tamil sentences that a young child can understand.
- Use everyday Tamil words. Avoid very complex or literary Tamil.
- Do not mix English words unless they are proper names.

IMPORTANT FOR CONTENT
- Read the student's question and understand the topic.
- Then explain ONLY that topic.
- For example, if the question is about an animal like a lion, explain:
  - what it is,
  - where it lives,
  - what it eats,
  - one or two simple facts.
- If the question is about maths like 'கூட்டல்', explain the maths idea in simple Tamil.

HARD RESTRICTIONS
- Do NOT talk about alphabet songs, letters, 'அகர வரிசை', 'எழுத்து', 'ஒலி', reading practice, or how to study UNLESS the question itself is about letters or reading.
- Do NOT invent classroom activities like 'let us sing the alphabet song', 'let us read the poem', 'circle the words', etc., unless the question explicitly asks.
- Do NOT ask the student to 'ask a question again'. Assume the sentence you see IS their question.
- Do NOT mention trees sharing fruits, children sharing fruits, poems about 'for' and 'on', or any similar reading-passage content unless the student's question itself mentions trees, fruits, or that poem.

HOW TO ANSWER
- First sentence: start explaining the topic directly.
- Use about 3 to 6 short sentences in total.
- Stay on-topic. Do not change the topic.

STUDENT QUESTION (in Tamil):
{question}

Now give your answer in Tamil, following ALL the rules above. Do not include any English. Do not include emojis or special symbols.
"""
            else:
                conversation_context = self.get_conversation_context()
                prompt = f"""You are a friendly, patient teacher for 1st and 2nd grade students.

GENERAL BEHAVIOUR
- Always answer clearly, using simple words and short sentences.
- Match the student's language choice through the rules given below.
- First, give a direct and clear answer to the student's question.
- Do not reply with vague questions like "What do you want to learn?" because the student has already asked something specific.
- Do not tell the student to "ask a question" again; assume the text you received already IS their question.
- Only after answering may you add one short suggestion or follow-up question, if it helps learning.
- Do not turn every answer into a quiz or worksheet unless the student clearly asks.

EXAMPLES (FOLLOW THIS STYLE)

Example 1:
Student (Tamil): "கூட்டல் பற்றி எனக்கு கற்றுக்கொடுங்கள்."
Teacher (Tamil): "கூட்டல் என்பது இரண்டு அல்லது அதற்கு மேற்பட்ட எண்களை ஒன்றாக சேர்த்து ஒரு புதிய எண்ணை பெறும் கணக்கு. உதாரணத்திற்கு, 2 மற்றும் 3 ஐ கூட்டினால் 5 கிடைக்கும். நாம் பழங்கள், பென்சில்கள் போன்ற பொருட்களை எண்ணி சேர்க்கும் போது கூட்டலைப் பயன்படுத்தலாம்."

Example 2:
Student (Tamil): "சிங்கங்களைப் பற்றி சொல்லுங்கள்."
Teacher (Tamil): "சிங்கம் ஒரு பெரிய மிருகம். அது பொதுவாக ஆப்பிரிக்கா புல்வெளிகளில் வாழ்கிறது. சிங்கம் இறைச்சி தின்று வாழ்கிறது. ஆண் சிங்கத்திற்கு பெரிய மயிர் வளையம் இருக்கும். பலர் சிங்கத்தை 'காட்டின் அரசன்' என்று அழைக்கிறார்கள்."

Example 3:
Student (English): "What is 2+2?"
Teacher (English): "2 + 2 = 4. This is an example of simple addition. When you put two objects together with two more objects, you get four objects in total."

Do NOT talk about trees sharing fruits, or poems about 'for' and 'on', unless the student's question itself mentions trees, fruits, or that poem.

{conversation_context}

STUDENT QUESTION: {question}

Provide a clear, direct, age-appropriate answer that follows the language and style rules below.
If the question is a simple factual or math question (for example 'what is 2+2' or 'கூட்டல் பற்றி'), give the correct answer and, if needed, a very short explanation.

Do NOT use emojis, symbols like *, +, =, :, or formatting characters such as markdown.
Respond only in clean, plain text suitable for text-to-speech systems.

{language_instruction}
"""

        if not self.llm:
            answer = "[ERROR] Language model is not available. Please check your GROQ_API_KEY in the .env file and restart the application."
        else:
            try:
                response = self.llm.invoke(prompt)
                answer = response.content.strip()
            except Exception as e:
                print(f"[WARNING] LLM Error: {e}")
                if "401" in str(e) or "invalid" in str(e).lower():
                    answer = "[ERROR] Invalid GROQ API key. Please check your GROQ_API_KEY in backend/.env file. Get a new key from the Groq console."
                elif not self.embeddings_available:
                    # Try to reload embeddings once per session
                    if not hasattr(self, '_embeddings_retry_attempted'):
                        self._embeddings_retry_attempted = True
                        print("[INFO] Attempting to reload embeddings...]")
                        if self.retry_embeddings_loading():
                            return self.query(question, top_k, target_language=normalized_language)

                    answer = "I'm running in offline mode right now. I can help with simple math problems like '5 + 3' or general conversations, but I cannot access educational documents. Check your internet connection and try restarting the application."
                else:
                    answer = f"I'm having trouble connecting to my language model right now: {e}. Please try again or check your API key."

        self.conversation_history.append({"role": "assistant", "content": answer})

        if len(self.conversation_history) > 24:
            self.conversation_history = self.conversation_history[-24:]

        return answer
